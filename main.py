from fastapi import (
    FastAPI,
    HTTPException,
    Depends,
    status,
    File,
    UploadFile,
    Form,
    Request,
)
from fastapi.staticfiles import StaticFiles
import os
import logging
import models  # models.py 全体をインポート
from routers import auth, users, chat, folders, upload, templates, images, video
from routers import memory as memory_router
from database import (
    engine,
    Base,
    get_db,
)  # SessionLocal はここでは直接使わないので削除、get_db を追加
from sqlalchemy.orm import Session  # SQLAlchemyのSession型をインポート
from sqlalchemy.sql import func  # SQLAlchemyのSQL関数(例: func.now())をインポート
from dependencies import (
    get_current_active_user,
)  # 認証済みユーザー取得用の依存関係をインポート
from models import (
    User,
    ChatSession,
    ChatMessage,
)  # ChatSession, ChatMessageも明示的にインポート
from dotenv import load_dotenv
from fastapi.middleware.cors import CORSMiddleware
from utils.openai_client import openai_client
from utils.logging_config import configure_logging
from anthropic import AsyncAnthropic
import google.generativeai as genai
from cohere import AsyncClient as AsyncCohereClient
from perplexipy import PerplexityClient  # 同期クライアントなので注意
import deepl
import boto3
from botocore.exceptions import NoCredentialsError, PartialCredentialsError
from fastapi.concurrency import run_in_threadpool  # 同期処理を非同期で実行するため
import base64
import fitz  # PyMuPDF
import uuid
from fastapi.responses import FileResponse  # ファイルダウンロード用
import aiofiles  # 非同期ファイル書き込み用

# オフィス文書処理用ライブラリ
from docx import Document as DocxDocument  # python-docx
from openpyxl import load_workbook  # openpyxl
from pptx import Presentation  # python-pptx
import json  # 今回の修正では直接使用していませんが、一般的に役立つため残します
import asyncio  # 今回の修正では直接使用していませんが、一般的に役立つため残します
from datetime import datetime, timezone

# ... (他のimport) ...
from pydantic import BaseModel, Field
from typing import Optional, Dict, Any, List
import schemas

# AI client functions are now imported from ai_clients.py
from ai_clients import (
    format_memories_for_prompt,
    get_openai_response,
    get_claude_response,
    get_cohere_response,
    get_perplexity_response,
    get_gemini_response,
    translate_with_deepl,
    FRIENDLY_TONE_SYSTEM_PROMPT
)
from ai_processing_flows import run_quality_chat_mode_flow, run_super_writing_orchestrator_flow
import shutil
import subprocess
import mimetypes
import traceback
from io import BytesIO

load_dotenv()
configure_logging()
logger = logging.getLogger(__name__)

# 生成されたファイルを保存するディレクトリを準備
GENERATED_FILES_DIR = "generated_files"
os.makedirs(GENERATED_FILES_DIR, exist_ok=True)

app = FastAPI()
app.mount("/static", StaticFiles(directory="static"), name="static")

@app.get("/", response_class=FileResponse)
async def read_index(request: Request):
    # Return the main UI file.
    # Ensure the path is correct relative to where main.py is executed.
    # Assuming 'static' directory is at the same level as main.py.
    return FileResponse("static/collaboration_ai_ui.html")

# models.Base.metadata.create_all(bind=engine)
# --- CORSミDLEWAREの設定 ---
origins = [
    "http://localhost",
    "http://localhost:3000",
    "http://localhost:8080",
    "http://127.0.0.1",
    "http://127.0.0.1:3000",
    "http://127.0.0.1:8080",
    "null",
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)
app.include_router(auth.router)
app.include_router(users.router)
app.include_router(chat.router)
app.include_router(folders.router)
app.include_router(upload.router)
app.include_router(templates.router)
app.include_router(images.router)
app.include_router(video.router)
app.include_router(memory_router.router)

# --- AIクライアントと設定の初期化 ---
# OpenAI (utils.openai_client で初期化済み)
app.state.openai_client = openai_client

# Anthropic
anthropic_api_key = os.getenv("ANTHROPIC_API_KEY")
if not anthropic_api_key:
    logger.warning(
        "警告: ANTHROPIC_API_KEY が設定されていません。Claude機能は利用できません。"
    )
    app.state.anthropic_client = None
else:
    app.state.anthropic_client = AsyncAnthropic(api_key=anthropic_api_key, timeout=600.0)

# Google Gemini
google_api_key = os.getenv("GOOGLE_API_KEY")
if not google_api_key:
    logger.warning(
        "警告: GOOGLE_API_KEY が設定されていません。Gemini機能は利用できません。"
    )
    app.state.gemini_vision_client = None
    app.state.gemini_pro_model = None
    app.state.gemini_flash_model = None
else:
    genai.configure(api_key=google_api_key)
    app.state.gemini_vision_client = genai.GenerativeModel(
        "gemini-2.5-pro-preview-05-06"
    )
    app.state.gemini_pro_model = genai.GenerativeModel(
        "gemini-2.5-pro-preview-05-06"
    )
    app.state.gemini_flash_model = genai.GenerativeModel(
        "gemini-2.5-flash-preview-04-17"
    )

# Cohere
cohere_api_key = os.getenv("COHERE_API_KEY")
if not cohere_api_key:
    logger.warning(
        "警告: COHERE_API_KEY が設定されていません。Cohere機能は利用できません。"
    )
    app.state.cohere_client = None
else:
    app.state.cohere_client = AsyncCohereClient(cohere_api_key, timeout=600)

# Perplexity
perplexity_api_key = os.getenv("PERPLEXITY_API_KEY")
if not perplexity_api_key:
    logger.warning(
        "警告: PERPLEXITY_API_KEY が設定されていません。Perplexity機能は利用できません。"
    )
    app.state.perplexity_sync_client = None
else:
    app.state.perplexity_sync_client = PerplexityClient(perplexity_api_key)

# DeepL
deepl_api_key = os.getenv("DEEPL_API_KEY")
if not deepl_api_key:
    logger.warning(
        "警告: DEEPL_API_KEY が設定されていません。翻訳機能は利用できません。"
    )
    app.state.deepl_translator = None
else:
    app.state.deepl_translator = deepl.Translator(deepl_api_key)

# AWS Textract
aws_region = os.getenv("AWS_DEFAULT_REGION") or os.getenv(
    "AWS_REGION", "ap-northeast-1"
)
try:
    app.state.textract_client = boto3.client("textract", region_name=aws_region)
    logger.info("AWS Textract client initialized for region: %s", aws_region)
except (NoCredentialsError, PartialCredentialsError) as e:
    logger.warning(
        "警告: AWS認証情報が見つからないか不完全です。Textract機能は利用できません。エラー: %s",
        e,
    )
    app.state.textract_client = None
except Exception as e:
    logger.warning(
        "警告: AWS Textractクライアントの初期化に失敗しました。エラー: %s", e
    )
    app.state.textract_client = None

# File processing helper functions
async def process_text_file(file: UploadFile) -> str:
    try:
        contents = await file.read()
        return contents.decode("utf-8", errors="replace")
    except Exception as e:
        logger.info(f"テキストファイルの読み込みエラー: {e}")
        raise
    finally:
        await file.seek(0)


async def process_image_with_vision_api(
    request: Request,
    file: UploadFile,
    original_prompt: str,
    client_identifier: str = "openai"
) -> str:
    if not request.app.state.openai_client:
        raise Exception("OpenAIクライアントが初期化されていません。")

    try:
        image_bytes = await file.read()
        await file.seek(0)

        base64_image = base64.b64encode(image_bytes).decode("utf-8")
        mime_type = file.content_type or "image/jpeg"

        user_image_prompt = "この画像の内容を詳細に説明してください。"
        if original_prompt:
            user_image_prompt += f"\nユーザーの主な関心事や、この画像をアップロードした背景にある質問は「{original_prompt}」です。これを踏まえて説明に関連性を持たせてください。"

        messages_for_api = [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": user_image_prompt},
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:{mime_type};base64,{base64_image}"},
                    },
                ],
            }
        ]
        # FRIENDLY_TONE_SYSTEM_PROMPT is imported from ai_clients.py
        if FRIENDLY_TONE_SYSTEM_PROMPT:
            messages_for_api.insert(
                0,
                {
                    "role": "system",
                    "content": FRIENDLY_TONE_SYSTEM_PROMPT
                    + "\n\nあなたはアップロードされた画像を分析し、ユーザーの関心事を考慮して内容を説明するAIです。",
                },
            )

        logger.info(
            f"OpenAI Vision API呼び出し準備 (モデル: gpt-4o): プロンプト='{user_image_prompt[:100]}...'"
        )
        api_response = await request.app.state.openai_client.chat.completions.create(
            model="gpt-4o",
            messages=messages_for_api,
            max_tokens=1000,
            temperature=0.5,
        )

        description = api_response.choices[0].message.content
        if not description or description.strip() == "":
            return f"画像「{file.filename}」の内容を認識できませんでした。"

        logger.info(
            f"OpenAI Vision APIからの応答 (冒頭): {description[:100].strip()}..."
        )
        return description.strip()

    except Exception as e:
        logger.info(f"OpenAI Vision APIでの画像処理中にエラー: {e}")
        traceback.print_exc()
        raise Exception(
            f"画像「{file.filename}」の処理中にエラーが発生しました: {str(e)}"
        )


async def process_pdf_with_ai(file: UploadFile) -> str:
    text_content = ""
    try:
        pdf_bytes = await file.read()
        await file.seek(0)

        def extract_text_from_pdf_sync(pdf_data: bytes) -> str:
            extracted_text = ""
            try:
                with fitz.open(stream=pdf_data, filetype="pdf") as doc:
                    for page_num in range(len(doc)):
                        page = doc.load_page(page_num)
                        extracted_text += page.get_text("text")
                        if page_num < len(doc) - 1:
                            extracted_text += "\n\n--- (次のページ) ---\n\n"
                return extracted_text
            except Exception as sync_e:
                logger.info(f"PyMuPDF (同期処理内) エラー: {sync_e}")
                raise Exception(
                    f"PDFからのテキスト抽出中にエラー (同期処理内): {str(sync_e)}"
                )

        logger.info(
            f"PyMuPDFによるPDF「{file.filename}」のテキスト抽出処理を開始します..."
        )
        text_content = await run_in_threadpool(extract_text_from_pdf_sync, pdf_bytes)

        if not text_content or text_content.strip() == "":
            return f"PDFファイル「{file.filename}」からテキストを抽出できませんでした。"

        logger.info(
            f"PDF「{file.filename}」からのテキスト抽出成功 (冒頭): {text_content[:200].strip()}..."
        )
        return text_content.strip()

    except Exception as e:
        logger.info(f"PDF「{file.filename}」の処理中にエラー: {e}")
        traceback.print_exc()
        raise Exception(
            f"PDFファイル「{file.filename}」の処理中にエラーが発生しました: {str(e)}"
        )


async def process_audio_with_speech_to_text(
    request: Request,
    file: UploadFile
) -> str:
    if not request.app.state.openai_client:
        raise Exception("OpenAIクライアントが初期化されていません。")

    try:
        audio_bytes = await file.read()
        await file.seek(0)

        logger.info(
            f"OpenAI Whisper API呼び出し準備 (モデル: whisper-1): ファイル名='{file.filename}', Content-Type='{file.content_type}'"
        )

        audio_file_for_api = BytesIO(audio_bytes)

        transcription_response = (
            await request.app.state.openai_client.audio.transcriptions.create(
                model="whisper-1",
                file=(file.filename, audio_file_for_api, file.content_type),
            )
        )

        transcribed_text = transcription_response.text

        if not transcribed_text or transcribed_text.strip() == "":
            return f"音声ファイル「{file.filename}」から文字を認識できませんでした。"

        logger.info(
            f"OpenAI Whisper APIからの文字起こし成功 (冒頭): {transcribed_text[:100].strip()}..."
        )
        return transcribed_text.strip()

    except Exception as e:
        logger.info(f"OpenAI Whisper APIでの音声処理中にエラー: {e}")
        traceback.print_exc()
        raise Exception(
            f"音声ファイル「{file.filename}」の処理中にエラーが発生しました: {str(e)}"
        )


async def process_docx_file(file_path: str) -> str:
    """DOCXファイルからテキストを抽出します。"""

    def extract_text_sync(path: str) -> str:
        try:
            doc = DocxDocument(path)
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            return "\n".join(full_text)
        except Exception as e:
            logger.info(f"DOCX処理エラー (同期処理内, path: {path}): {e}")
            raise Exception(f"DOCXファイルからのテキスト抽出中にエラー: {str(e)}")

    logger.info(
        f"python-docxによるDOCXファイル「{os.path.basename(file_path)}」のテキスト抽出処理を開始します..."
    )
    text_content = await run_in_threadpool(extract_text_sync, file_path)
    if not text_content or text_content.strip() == "":
        return f"DOCXファイル「{os.path.basename(file_path)}」からテキストを抽出できませんでした。"
    logger.info(
        f"DOCX「{os.path.basename(file_path)}」からのテキスト抽出成功 (冒頭): {text_content[:200].strip()}..."
    )
    return text_content.strip()


async def process_xlsx_file(file_path: str) -> str:
    """XLSXファイルから各セルのテキスト情報を抽出します。"""

    def extract_text_sync(path: str) -> str:
        try:
            workbook = load_workbook(filename=path, read_only=True, data_only=True)
            full_text = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                full_text.append(f"\n--- シート: {sheet_name} ---\n")
                for row in sheet.iter_rows():
                    row_text = []
                    for cell in row:
                        if cell.value is not None:
                            row_text.append(str(cell.value))
                    if row_text:
                        full_text.append("\t".join(row_text))
            return "\n".join(full_text)
        except Exception as e:
            logger.info(f"XLSX処理エラー (同期処理内, path: {path}): {e}")
            raise Exception(f"XLSXファイルからのテキスト抽出中にエラー: {str(e)}")

    logger.info(
        f"openpyxlによるXLSXファイル「{os.path.basename(file_path)}」のテキスト抽出処理を開始します..."
    )
    text_content = await run_in_threadpool(extract_text_sync, file_path)
    if not text_content or text_content.strip() == "":
        return f"XLSXファイル「{os.path.basename(file_path)}」からテキストを抽出できませんでした。"
    logger.info(
        f"XLSX「{os.path.basename(file_path)}」からのテキスト抽出成功 (冒頭): {text_content[:200].strip()}..."
    )
    return text_content.strip()


async def process_pptx_file(file_path: str) -> str:
    """PPTXファイルから各スライドのテキスト情報を抽出します。"""

    def extract_text_sync(path: str) -> str:
        try:
            presentation = Presentation(path)
            full_text = []
            for i, slide in enumerate(presentation.slides):
                full_text.append(f"\n--- スライド {i + 1} ---\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                full_text.append(run.text)
                    elif hasattr(shape, "text") and shape.text:
                        full_text.append(shape.text)
            return "\n".join(full_text)
        except Exception as e:
            logger.info(f"PPTX処理エラー (同期処理内, path: {path}): {e}")
            raise Exception(f"PPTXファイルからのテキスト抽出中にエラー: {str(e)}")

    logger.info(
        f"python-pptxによるPPTXファイル「{os.path.basename(file_path)}」のテキスト抽出処理を開始します..."
    )
    text_content = await run_in_threadpool(extract_text_sync, file_path)
    if not text_content or text_content.strip() == "":
        return f"PPTXファイル「{os.path.basename(file_path)}」からテキストを抽出できませんでした。"
    logger.info(
        f"PPTX「{os.path.basename(file_path)}」からのテキスト抽出成功 (冒頭): {text_content[:200].strip()}..."
    )
    return text_content.strip()


async def convert_markdown_to_format_with_pandoc(
    markdown_content: str,
    output_filename_base: str,
    output_format: str,
    temp_dir: str = GENERATED_FILES_DIR,
) -> Optional[str]:
    """Pandocを使用してMarkdownを指定された形式に変換し、一時ファイルに保存する。"""
    input_md_path = os.path.join(temp_dir, f"{output_filename_base}_temp.md")
    output_file_path = os.path.join(temp_dir, f"{output_filename_base}.{output_format}")

    try:
        async with aiofiles.open(input_md_path, "w", encoding="utf-8") as md_file:
            await md_file.write(markdown_content)

        PANDOC_EXECUTABLE_PATH = (
            r"C:\Program Files\Pandoc\pandoc.exe"
        )

        pandoc_cmd = [
            PANDOC_EXECUTABLE_PATH,
            input_md_path,
            "-o",
            output_file_path,
        ]
        if output_format == "pdf":
            pandoc_cmd.extend(
                [
                    "--pdf-engine=lualatex",
                    "-V",
                    "documentclass=ltjarticle",
                    "-V",
                    "mainfont=IPAexMincho",
                    "-V",
                    "sansfont=IPAexGothic",
                    "-V",
                    "monofont=IPAexGothic",
                ]
            )
        elif output_format == "docx":
            pass

        logger.info(f"Pandocコマンド実行準備: {' '.join(pandoc_cmd)}")

        def run_pandoc_sync():
            process = subprocess.run(
                pandoc_cmd, capture_output=True, text=True, check=False
            )
            if process.returncode != 0:
                error_message = f"Pandoc変換エラー (フォーマット: {output_format}):\nSTDOUT:\n{process.stdout}\nSTDERR:\n{process.stderr}"
                logger.info(error_message)
                raise Exception(
                    f"Pandoc failed with exit code {process.returncode}. STDERR: {process.stderr[:500]}"
                )
            return True

        await run_in_threadpool(run_pandoc_sync)

        logger.info(f"Pandocによるファイル「{output_file_path}」の生成に成功しました。")
        return output_file_path

    except Exception as e:
        logger.info(
            f"Pandoc変換処理中にエラー ({output_format} へ変換しようとしていました): {e}"
        )
        traceback.print_exc()
        if os.path.exists(output_file_path):
            try:
                os.remove(output_file_path)
            except OSError:
                pass
        return None
    finally:
        if os.path.exists(input_md_path):
            try:
                os.remove(input_md_path)
            except OSError:
                pass

@app.post("/translate", response_model=schemas.TranslationResponse)
async def translate_endpoint(
    translation_payload: schemas.TranslationRequest,
    fastapi_request: Request,
    current_user: models.User = Depends(get_current_active_user),
):
    translated = await translate_with_deepl(fastapi_request, translation_payload.text, translation_payload.target_lang)
    return {"translated_text": translated}


@app.get("/download_generated_file/{filename}")
async def download_generated_file(
    filename: str, current_user: models.User = Depends(get_current_active_user)
):
    if ".." in filename or filename.startswith("/"):
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST, detail="無効なファイル名です。"
        )

    file_path = os.path.join(GENERATED_FILES_DIR, filename)

    if not os.path.exists(file_path) or not os.path.isfile(file_path):
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail="ファイルが見つかりません。"
        )
    return FileResponse(
        path=file_path, filename=filename, media_type="application/octet-stream"
    )

@app.post("/collaborative_answer_v2", response_model=schemas.CollaborativeResponseV2)
async def collaborative_answer_mode_endpoint(
    fastapi_request: Request,
    prompt: str = Form(...),
    mode: str = Form(...),
    session_id: Optional[int] = Form(None),
    char_count: Optional[int] = Form(None),
    genre: Optional[str] = Form(None),
    file: Optional[UploadFile] = File(None),
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_active_user),
) -> schemas.CollaborativeResponseV2:
    original_prompt_from_user = prompt.strip()
    current_mode = mode.lower()
    current_session_id_from_request = session_id
    desired_char_count = char_count

    logger.info("--- collaborative_answer_mode_endpoint ---")
    logger.info(
        f"Received mode from frontend: '{mode}' (raw), Processed current_mode: '{current_mode}'"
    )

    user_memories_from_request: Optional[List[schemas.UserMemoryResponse]] = None
    try:
        user_memories_from_request = (
            db.query(models.UserMemory)
            .filter(models.UserMemory.user_id == current_user.id)
            .order_by(
                models.UserMemory.priority.desc(),
                models.UserMemory.updated_at.desc(),
            )
            .all()
        )
    except Exception as e:
        logger.info(f"ユーザーメモリ取得中にエラー: {e}")
        user_memories_from_request = None

    logger.info(
        f"\nリクエスト受信: UserID={current_user.id}, SessionID(Req)={current_session_id_from_request}, Prompt='{original_prompt_from_user[:50].strip()}...', Mode='{current_mode}', File: {file.filename if file else 'なし'}"
    )

    response_shell = schemas.CollaborativeResponseV2(
        prompt=original_prompt_from_user,
        mode_executed=current_mode,
        processed_session_id=current_session_id_from_request,
    )

    final_prompt_for_ai_flow = original_prompt_from_user
    processed_file_text_for_ai: Optional[str] = None

    if file and file.filename:
        try:
            MAX_FILE_SIZE_MB = 50
            MAX_FILE_SIZE_BYTES = MAX_FILE_SIZE_MB * 1024 * 1024
            temp_file_path = None
            original_filename = file.filename

            base_temp_dir = "temp_uploads"
            os.makedirs(base_temp_dir, exist_ok=True)
            file_ext = os.path.splitext(original_filename)[1]
            safe_filename = f"{uuid.uuid4()}{file_ext}"
            temp_file_path = os.path.join(base_temp_dir, safe_filename)

            current_size = 0
            with open(temp_file_path, "wb") as f:
                while chunk := await file.read(1024 * 1024):
                    current_size += len(chunk)
                    if current_size > MAX_FILE_SIZE_BYTES:
                        f.close()
                        os.remove(temp_file_path)
                        raise HTTPException(
                            status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE,
                            detail=f"ファイル「{original_filename}」のサイズが大きすぎます (最大{MAX_FILE_SIZE_MB}MB)。",
                        )
                    f.write(chunk)
            await file.seek(0)

            mime_type = file.content_type or mimetypes.guess_type(original_filename)[0]
            logger.info(f"ファイル処理開始: {original_filename}, MIMEタイプ: {mime_type}")

            file_processing_log: Optional[schemas.IndividualAIResponse] = None

            if mime_type:
                if mime_type.startswith("text/") or any(original_filename.lower().endswith(ext) for ext in ['.py', '.js', '.html', '.css', '.md']):
                    processed_file_text_for_ai = await process_text_file(file)
                    file_processing_log = schemas.IndividualAIResponse(source="ファイル処理(テキスト)", response=f"テキスト系ファイル「{original_filename}」の内容を読み込みました。")
                elif mime_type.startswith("image/"):
                    processed_file_text_for_ai = await process_image_with_vision_api(fastapi_request, file, original_prompt_from_user)
                    file_processing_log = schemas.IndividualAIResponse(source="ファイル処理(画像認識)", response=f"画像ファイル「{original_filename}」の内容をAIが認識しました。")
                elif mime_type == "application/pdf":
                    processed_file_text_for_ai = await process_pdf_with_ai(file)
                    file_processing_log = schemas.IndividualAIResponse(source="ファイル処理(PDF)", response=f"PDFファイル「{original_filename}」からテキスト情報を抽出しました。")
                elif mime_type.startswith("audio/") or any(original_filename.lower().endswith(ext) for ext in ['.mp3', '.mp4', '.m4a', '.wav', '.webm', '.mpeg']):
                    processed_file_text_for_ai = await process_audio_with_speech_to_text(fastapi_request, file)
                    file_processing_log = schemas.IndividualAIResponse(source="ファイル処理(音声認識)", response=f"音声ファイル「{original_filename}」を文字起こししました。")
                elif original_filename.lower().endswith(".docx"):
                    processed_file_text_for_ai = await process_docx_file(temp_file_path)
                    file_processing_log = schemas.IndividualAIResponse(source="ファイル処理(DOCX)", response=f"DOCXファイル「{original_filename}」からテキスト情報を抽出しました。")
                elif original_filename.lower().endswith(".xlsx"):
                    processed_file_text_for_ai = await process_xlsx_file(temp_file_path)
                    file_processing_log = schemas.IndividualAIResponse(source="ファイル処理(XLSX)", response=f"XLSXファイル「{original_filename}」からテキスト情報を抽出しました。")
                elif original_filename.lower().endswith(".pptx"):
                    processed_file_text_for_ai = await process_pptx_file(temp_file_path)
                    file_processing_log = schemas.IndividualAIResponse(source="ファイル処理(PPTX)", response=f"PPTXファイル「{original_filename}」からテキスト情報を抽出しました。")
                else:
                    file_processing_log = schemas.IndividualAIResponse(source="ファイル処理", error=f"ファイル形式「{mime_type or '不明'}」は現在直接処理できません。ファイル名「{original_filename}」")
            else:
                 file_processing_log = schemas.IndividualAIResponse(source="ファイル処理", error=f"ファイル「{original_filename}」の形式を特定できませんでした。")

            response_shell.file_processing_step = file_processing_log
            if file_processing_log and file_processing_log.error:
                final_prompt_for_ai_flow = (
                    f"ユーザーはファイル「{original_filename}」をアップロードしましたが、その内容の処理中にエラーが発生したため利用できません。\n"
                    f"ファイル内容に関する通知: 「{file_processing_log.error}」\n"
                    f"この状況を踏まえ、以下のユーザーの指示に対応してください。\n---\n"
                    f"ユーザーの指示: 「{original_prompt_from_user}」"
                )
            elif processed_file_text_for_ai:
                final_prompt_for_ai_flow = (
                    f"ユーザーは次のファイルをアップロードしました。\n"
                    f"ファイル名: {original_filename}\n"
                    f"ファイルから抽出・認識された内容の要約または全文:\n```text\n{processed_file_text_for_ai[:2000].strip()}...\n```\n"
                    f"(上記はアップロードされたファイルの内容です。これを踏まえて、以下のユーザーの指示に対応してください。)\n---\n"
                    f"ユーザーの指示: 「{original_prompt_from_user}」"
                )
            else:
                final_prompt_for_ai_flow = (
                    f"ユーザーはファイル「{original_filename}」をアップロードしましたが、抽出可能なテキスト情報が含まれていませんでした。\n"
                    f"この状況を踏まえ、以下のユーザーの指示に対応してください。\n---\n"
                    f"ユーザーの指示: 「{original_prompt_from_user}」"
                )
            response_shell.prompt = final_prompt_for_ai_flow

        except HTTPException as he:
            logger.info(f"ファイル処理HTTPエラー: {he.detail}")
            response_shell.overall_error = he.detail
            return response_shell
        except Exception as e:
            logger.info(f"ファイル処理中に予期せぬエラー: {e}\n{traceback.format_exc()}")
            response_shell.overall_error = f"ファイルの処理中に予期せぬエラーが発生しました: {str(e)}"
            return response_shell
        finally:
            if file: await file.close()
            if temp_file_path and os.path.exists(temp_file_path):
                try: os.remove(temp_file_path)
                except OSError: pass

    active_session: Optional[models.ChatSession] = None
    initial_user_prompt_for_session: Optional[str] = None
    chat_history_for_ai: List[Dict[str, str]] = []

    if current_session_id_from_request:
        active_session = db.query(models.ChatSession).filter(models.ChatSession.id == current_session_id_from_request, models.ChatSession.user_id == current_user.id).first()
        if not active_session:
            raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="指定されたチャットセッションが見つからないか、アクセス権がありません。")
        response_shell.processed_session_id = active_session.id
        first_user_message_db = db.query(models.ChatMessage).filter(models.ChatMessage.chat_session_id == active_session.id, models.ChatMessage.role == "user").order_by(models.ChatMessage.created_at.asc()).first()
        initial_user_prompt_for_session = first_user_message_db.content if first_user_message_db else original_prompt_from_user
        past_messages_db = db.query(models.ChatMessage).filter(models.ChatMessage.chat_session_id == active_session.id).order_by(models.ChatMessage.created_at.asc()).all()
        for msg_db in past_messages_db:
            chat_history_for_ai.append({"role": "assistant" if msg_db.role == "ai" else msg_db.role, "content": msg_db.content})
    else:
        potential_title = original_prompt_from_user.splitlines()[0][:50].strip() or "新しいチャット"
        active_session = models.ChatSession(user_id=current_user.id, title=potential_title, status="loading", mode=current_mode)
        db.add(active_session)
        db.commit()
        db.refresh(active_session)
        response_shell.processed_session_id = active_session.id
        initial_user_prompt_for_session = original_prompt_from_user

    user_message_db = models.ChatMessage(chat_session_id=active_session.id, user_id=current_user.id, role="user", content=original_prompt_from_user)
    db.add(user_message_db)
    active_session.updated_at = func.now()
    active_session.status = "loading"
    db.commit()
    db.refresh(user_message_db)
    db.refresh(active_session)

    try:
        if current_mode == "balance":
            response_shell = await run_balance_mode_flow(
                original_prompt=final_prompt_for_ai_flow,
                response_shell=response_shell,
                chat_history_for_ai=list(chat_history_for_ai),
                initial_user_prompt_for_session=initial_user_prompt_for_session,
                request=fastapi_request,
                user_memories=user_memories_from_request,
            )
        elif current_mode in ("search6", "supersearch"):
            response_shell = await run_super_search_mode_flow(
                original_prompt=final_prompt_for_ai_flow,
                response_shell=response_shell,
                chat_history_for_ai=list(chat_history_for_ai),
                initial_user_prompt_for_session=initial_user_prompt_for_session,
                user_memories=user_memories_from_request,
                request=fastapi_request,
            )
        elif current_mode == "superwriting":
            response_shell = await run_super_writing_orchestrator_flow(
                original_prompt=final_prompt_for_ai_flow,
                genre=genre,
                response_shell=response_shell,
                chat_history_for_ai=list(chat_history_for_ai),
                initial_user_prompt_for_session=initial_user_prompt_for_session,
                user_memories=user_memories_from_request,
                desired_char_count=desired_char_count,
                request=fastapi_request,
            )
        elif current_mode == "fastchat":
            response_shell = await run_fast_chat_mode_flow(
                original_prompt=final_prompt_for_ai_flow,
                response_shell=response_shell,
                chat_history_for_ai=list(chat_history_for_ai),
                initial_user_prompt_for_session=initial_user_prompt_for_session,
                user_memories=user_memories_from_request,
                request=fastapi_request,
            )
        else:
            raise HTTPException(
                status_code=400,
                detail=f"無効なモード「{current_mode}」が指定されました。",
            )

        final_ai_response_content_for_db: str = ""
        final_ai_source_text_for_db: str = f"Unknown AI ({mode} mode)"

        if (
            response_shell.step7_final_answer_v2_openai
            and response_shell.step7_final_answer_v2_openai.response
        ):
            final_ai_response_content_for_db = response_shell.step7_final_answer_v2_openai.response
            final_ai_source_text_for_db = (
                response_shell.step7_final_answer_v2_openai.source
                or f"Final Output ({mode.capitalize()} Mode)"
            )
        elif response_shell.overall_error:
            final_ai_response_content_for_db = (
                f"処理中にエラーが発生しました: {response_shell.overall_error}"
            )
            final_ai_source_text_for_db = f"Error in {mode.capitalize()} Mode"
        else:
            final_ai_response_content_for_db = "AIから有効な応答がありませんでした。"
            final_ai_source_text_for_db = (
                f"No Valid Response in {mode.capitalize()} Mode"
            )

        if final_ai_response_content_for_db and active_session and active_session.id:
            ai_message_db = models.ChatMessage(
                chat_session_id=active_session.id,
                role="ai",
                content=final_ai_response_content_for_db,
                ai_model=final_ai_source_text_for_db,
                user_id=None,
            )
            db.add(ai_message_db)
            active_session.updated_at = func.now()
            active_session.status = "complete"
            db.add(active_session)
            db.commit()
            db.refresh(ai_message_db)
            db.refresh(active_session)
            logger.info(
                f"AIレスポンス保存成功: MsgID={ai_message_db.id}, SessionID={active_session.id}, Source='{final_ai_source_text_for_db}'"
            )
        elif response_shell.overall_error:
            logger.info(
                f"AI処理でエラー発生のためDBへのAI応答保存をスキップ: {response_shell.overall_error}"
            )
            if active_session:
                active_session.status = "error"
                db.commit()
                db.refresh(active_session)
        else:
            logger.info(
                f"AIからの最終応答が見つからないか内容が空のためDBへのAI応答保存をスキップ: Mode='{mode}'"
            )
            if active_session:
                active_session.status = "complete_no_response"
                db.commit()
                db.refresh(active_session)

    except ValueError as ve:
        error_message = f"モード '{mode}' の処理中にエラーが発生しました: {str(ve)}"
        logger.info(
            f"ValueError in collaborative_answer_mode_endpoint: {error_message}"
        )
        response_shell.overall_error = error_message
        return response_shell
    except HTTPException as he:
        raise he
    except Exception as e:
        error_trace = traceback.format_exc()
        error_message = f"予期せぬエラーが発生しました: {str(e)}"
        logger.info(
            f"Unexpected Error in collaborative_answer_mode_endpoint: {error_message}\nTrace: {error_trace}"
        )
        response_shell.overall_error = f"サーバー処理中にエラーが発生しました: {str(e)}"
        return response_shell

    output_format_match = re.search(
        r"「(.+?)形式で出力」|「(.+?)として出力」|出力形式は\s*([a-zA-Z0-9]+)",
        original_prompt_from_user,
        re.IGNORECASE,
    )
    requested_output_format_str: Optional[str] = None
    if output_format_match:
        for group in output_format_match.groups():
            if group:
                requested_output_format_str = group.strip().lower()
                break

    if (
        requested_output_format_str
        and response_shell.step7_final_answer_v2_openai
        and response_shell.step7_final_answer_v2_openai.response
    ):
        final_ai_text_content_for_file = (
            response_shell.step7_final_answer_v2_openai.response
        )
        supported_text_extensions = [
            "txt", "md", "markdown", "json", "py", "html", "css", "js", "csv", "xml", "yaml", "yml", "log",
        ]
        actual_extension = requested_output_format_str
        if requested_output_format_str == "markdown": actual_extension = "md"
        elif requested_output_format_str == "python": actual_extension = "py"
        elif requested_output_format_str == "javascript": actual_extension = "js"

        output_file_generated = False
        unsupported_output_message = f"\n\n（システムより追記: ご指定の「{requested_output_format_str}」形式でのファイル出力は現在サポートされていないか、処理に失敗しました。テキストでの回答となります。）"

        if actual_extension in supported_text_extensions:
            try:
                session_prefix = f"s{active_session.id}_" if active_session and active_session.id else ""
                user_prefix = f"u{current_user.id}_"
                unique_id = str(uuid.uuid4())[:8]
                base_filename = f"{user_prefix}{session_prefix}output_{unique_id}"
                output_filename_with_ext = f"{base_filename}.{actual_extension}"
                full_output_path = os.path.join(GENERATED_FILES_DIR, output_filename_with_ext)
                async with aiofiles.open(full_output_path, "w", encoding="utf-8") as f:
                    await f.write(final_ai_text_content_for_file)
                response_shell.generated_download_url = f"/download_generated_file/{output_filename_with_ext}"
                response_shell.generated_file_name = output_filename_with_ext
                output_file_generated = True
                logger.info(f"テキストファイル生成成功: {full_output_path}")
            except Exception as e_file_write:
                logger.info(f"ファイル「{output_filename_with_ext if 'output_filename_with_ext' in locals() else 'unknown'}」の書き出しに失敗: {e_file_write}")
                traceback.print_exc()
        elif actual_extension == "pdf" or actual_extension == "docx":
            try:
                session_prefix = f"s{active_session.id}_" if active_session and active_session.id else ""
                user_prefix = f"u{current_user.id}_"
                unique_id = str(uuid.uuid4())[:8]
                base_filename_for_conversion = f"{user_prefix}{session_prefix}output_{unique_id}"
                generated_file_full_path = await convert_markdown_to_format_with_pandoc(
                    markdown_content=final_ai_text_content_for_file,
                    output_filename_base=base_filename_for_conversion,
                    output_format=actual_extension,
                    temp_dir=GENERATED_FILES_DIR,
                )
                if generated_file_full_path and os.path.exists(generated_file_full_path):
                    generated_filename_only = os.path.basename(generated_file_full_path)
                    response_shell.generated_download_url = f"/download_generated_file/{generated_filename_only}"
                    response_shell.generated_file_name = generated_filename_only
                    output_file_generated = True
                    logger.info(f"{actual_extension.upper()}ファイル生成成功: {generated_file_full_path}")
                else:
                    logger.info(f"Pandocによる{actual_extension.upper()}ファイル生成に失敗しました。")
            except Exception as e_conversion:
                logger.info(f"{actual_extension.upper()}への変換処理中にエラー: {e_conversion}")
                traceback.print_exc()

        if not output_file_generated and response_shell.step7_final_answer_v2_openai:
            response_shell.step7_final_answer_v2_openai.response = (response_shell.step7_final_answer_v2_openai.response or "") + unsupported_output_message

    logger.info(f"Endpoint (normal path) is about to return response_shell.")
    if response_shell:
        logger.info(f"Final content of response_shell: {response_shell.model_dump_json(indent=2)}")
    else:
        logger.info("Error: response_shell is None before returning from endpoint.")
        raise HTTPException(status_code=500, detail="サーバー内部エラー: レスポンスオブジェクトがnullです。")
    return response_shell

async def run_balance_mode_flow(
    original_prompt: str,
    response_shell: schemas.CollaborativeResponseV2,
    chat_history_for_ai: List[Dict[str, str]],
    initial_user_prompt_for_session: Optional[str],
    request: Request,
    user_memories: Optional[List[schemas.UserMemoryResponse]] = None,
) -> schemas.CollaborativeResponseV2:
    """High quality chat mode flow (Perplexity -> Claude)."""
    logger.info("\n--- ハイクオリティモード開始 (簡易) ---")
    res = await run_quality_chat_mode_flow(
        original_prompt=original_prompt,
        response_shell=response_shell,
        chat_history_for_ai=chat_history_for_ai,
        initial_user_prompt_for_session=initial_user_prompt_for_session,
        user_memories=user_memories,
        request=request,
    )
    logger.info("--- ハイクオリティモード終了 ---")
    logger.info(
        f"run_balance_mode_flow が返却する response_shell の内容 (JSON): {res.model_dump_json(indent=2) if res else None}"
    )
    return res


async def run_super_search_mode_flow(
    original_prompt: str,
    response_shell: schemas.CollaborativeResponseV2,
    chat_history_for_ai: List[Dict[str, str]],
    initial_user_prompt_for_session: Optional[str],
    request: Request,
    user_memories: Optional[List[schemas.UserMemoryResponse]] = None,
) -> schemas.CollaborativeResponseV2:
    """Five-step Perplexity search with final Claude summary."""
    perplexity_client = request.app.state.perplexity_sync_client
    claude_client = request.app.state.anthropic_client
    if not perplexity_client or not claude_client:
        response_shell.step7_final_answer_v2_openai = schemas.IndividualAIResponse(
            source="SuperSearch",
            error="必要なAIクライアントが初期化されていません。",
        )
        response_shell.overall_error = "必要なAIクライアントが初期化されていません。"
        return response_shell

    from datetime import datetime
    current_dt = datetime.utcnow().strftime("%Y-%m-%d %H:%M:%S")

    def dedup_lines(text: str) -> str:
        seen = set()
        lines = []
        for line in text.splitlines():
            stripped = line.strip()
            if stripped and stripped not in seen:
                lines.append(line)
                seen.add(stripped)
        return "\n".join(lines)

    def enhance_output_formatting(text: str) -> str:
        import re as _re
        return _re.sub(r"\n(?=\S)", "\n\n", text)

    results: List[str] = []
    summaries: List[str] = []
    all_perplexity_step_results: List[schemas.IndividualAIResponse] = []

    for step in range(5):
        if step == 0:
            prompt = (
                f"{current_dt} 時点での最新情報を調査してください。以下のテーマについて、1000文字以上の密度ある回答を作成してください：\n\n"
                f"「{original_prompt}」\n\n"
                "この出力を作成する前に、必ず以下を厳守してください：\n"
                "1. ユーザーメモリ（User Memory）を全て参照し、重要な指示・要望があれば反映してください（ただし、プロンプト内容を最優先）。\n"
                "2. それまでのチャット履歴全体を読み込み、会話の流れや一貫性、話の背景、過去に触れた話題を確認してください。\n"
                "3. トークン制限は無視して構いません。1000文字未満の出力は認められません。"
            )
        else:
            summary_lines = "\n".join(
                [f"- 要約{i+1}：{summaries[i][:300]}..." for i in range(step)]
            )
            prompt = (
                f"{current_dt} 時点での最新情報を再調査してください。\n\n"
                f"対象テーマ：「{original_prompt}」\n\n"
                "前回までに得られた情報（重複禁止）：\n" + summary_lines + "\n\n"
                "今回のタスク：\n"
                "- 上記の要約と重複しない新しい観点・角度・情報源から、最低1000文字以上の新たな回答を出力してください。\n"
                "- 視点が異なる専門領域、実例、地理的差異、時系列変化、統計、法規制、社会的論争などからアプローチして構いません。\n"
                "- 単なる言い換えや抽象化ではなく、具体的な追加情報や文脈の広がりを含めてください。\n\n"
                "必ず以下を実施：\n"
                "1. ユーザーメモリ（User Memory）全体を読み取り、使えそうな情報があれば反映。\n"
                "2. チャット履歴を全読破し、文脈や目的意識を保ち、一貫性のある内容に。\n"
                "3. 出力は必ず1000文字以上。満たさない場合は再度新しい観点で再検索してください。"
            )
        res = await get_perplexity_response(
            request=request,
            prompt_for_perplexity=prompt,
            model="sonar-reasoning-pro",
            user_memories=user_memories,
            initial_user_prompt=initial_user_prompt_for_session,
        )
        all_perplexity_step_results.append(res)
        text = res.response or ""
        text = dedup_lines(text)
        tries = 0
        while len(text) < 1000 and not res.error and tries < 2:
            extra = await get_perplexity_response(
                request=request,
                prompt_for_perplexity=(
                    "前回の情報と重複しない新しい観点から、さらに詳しく1000文字以上で説明してください:\n"
                    + prompt
                ),
                model="sonar-reasoning-pro",
                user_memories=user_memories,
                initial_user_prompt=initial_user_prompt_for_session,
            )
            if extra.response:
                text += "\n" + extra.response
                text = dedup_lines(text)
            tries += 1
        results.append(text)
        summaries.append(text[:500])

    summary_prompt = (
        "以下に、5段階に分けて収集した情報があります。\nこれらを1つの統一された読みやすく魅力的な日本語文章にまとめてください。\n\n"
        "【絶対条件】\n"
        "- 出力は5000文字以上でなければなりません。\n"
        "- 内容は削除せず、むしろ適度に補足しながら流れを整えてください。\n"
        "- 読者の興味を引く構成（導入→展開→結論）で、話題の全体像を掴めるようにしてください。\n"
        "- 論点が散らばらないよう、全体構成を練ってください。\n\n"
        "【使用する情報（5つの視点）】\n---\n①：" + results[0] + "\n---\n②：" + results[1] + "\n---\n③：" + results[2] + "\n---\n④：" + results[3] + "\n---\n⑤：" + results[4] + "\n---\n\n"
        "【開始前に必ず実行すること】\n"
        "1. ユーザーメモリ（User Memory）全体を精読し、内容に反映可能な情報があれば自然な形で取り入れてください。\n"
        "2. チャット履歴全体を全読破し、文脈、トーン、一貫性、以前の議論と矛盾しないよう注意してください。\n"
        "3. トークン制限や処理時間は一切考慮せず、最高の内容品質を最優先としてください。"
    )
    claude_system_instruction = "あなたは、複数の検索結果を統合し、詳細かつ網羅的な一つの日本語の報告書にまとめる専門家です。提供された情報を元に、指示通り5000文字以上で、内容を削らず、構成を整えて出力してください。"
    summary_res = await get_claude_response(
        request=request,
        prompt_text=summary_prompt,
        system_instruction=claude_system_instruction,
        model="claude-opus-4-20250514",
        chat_history=chat_history_for_ai,
        initial_user_prompt=initial_user_prompt_for_session,
        user_memories=user_memories,
    )
    formatted_response = (enhance_output_formatting(summary_res.response) if summary_res.response else summary_res.response)
    response_shell.step7_final_answer_v2_openai = schemas.IndividualAIResponse(
        source="Claude (claude-opus-4-20250514)", response=formatted_response, error=summary_res.error,
    )
    # response_shell.search_fragments = all_perplexity_step_results # Prevent search fragments from being sent
    response_shell.search_fragments = []
    return response_shell

async def run_writing_mode_flow(
    original_prompt: str,
    response_shell: schemas.CollaborativeResponseV2,
    chat_history_for_ai: List[Dict[str, str]],
    initial_user_prompt_for_session: Optional[str],
    request: Request,
    user_memories: Optional[List[schemas.UserMemoryResponse]] = None,
) -> schemas.CollaborativeResponseV2:
    logger.info("\n--- 執筆特化モード開始 ---")
    # Ensure `request` is used for AI helper calls within this function.
    steps_executed: List[schemas.IndividualAIResponse] = []
    # ... (variable initializations remain the same) ...
    defined_requirements = ""
    finalized_requirements = ""
    article_outline = ""
    draft_content = ""
    review_and_suggestions = ""
    revised_draft_content = ""
    final_article_content = ""

    current_chat_history_for_this_turn = list(chat_history_for_ai)
    current_chat_history_for_this_turn.append({"role": "user", "content": original_prompt})
    logger.info(f"Writing Mode: このターンでAIに渡す完全な履歴は {len(current_chat_history_for_this_turn)} 件, メモリ: {len(user_memories) if user_memories else 0}件")

    try:
        logger.info("執筆モード ステップW0: 要件確認とテーマ深掘り中...")
        w0_system_instruction = (
            "あなたはユーザーの多様な執筆リクエストを分析し、高品質なコンテンツを生成するために必要な構成要素を明確化する専門の編集アシスタントです。\n"
            f"この会話全体の主要な目的は「{initial_user_prompt_for_session}」です。\n"
        )
        w0_user_prompt = (f"ユーザーは以下の内容の執筆を希望しています。\nユーザーの執筆リクエスト: 「{original_prompt}」\n\n")
        w0_res = await get_gemini_response(
            request=request, prompt_text=w0_user_prompt, system_instruction=w0_system_instruction, model_name="gemini-2.5-pro-preview-05-06",
            chat_history=list(current_chat_history_for_this_turn), initial_user_prompt=initial_user_prompt_for_session, user_memories=user_memories,
        )
        steps_executed.append(w0_res)
        if w0_res.error or not w0_res.response: raise ValueError(f"執筆モード ステップW0 (要件確認) 失敗: {w0_res.error or '応答がありませんでした。'}")
        defined_requirements = w0_res.response
        logger.info(f"ステップW0 - 定義された執筆要件 (冒頭):\n{defined_requirements[:300].strip()}...")

        supplemented_info = f"""
【開発者による補足情報（ステップW0のAIからの質問への一般的な回答方針）】
- 執筆物の種類: ユーザーの最初の執筆リクエスト「{original_prompt}」から最大限推測してください。不明確な場合は、最も可能性の高い種類を仮定するか、複数の可能性を提示してユーザーに選択を促すような形で進めてください。
- ターゲット読者: 「{original_prompt}」から推測される読者層、または執筆物の種類から一般的に想定される読者層を考慮してください。
- トーン＆マナー: 「{original_prompt}」に含まれる雰囲気（例：「面白おかしく」「真剣に」「感動的に」など）を最優先してください。
- 主要な要素: 「{original_prompt}」に含まれるキーワード、テーマ、キャラクター、出来事などを中心に据えてください。
- ボリューム: ステップW0のAIがユーザーの指示や内容から提案したボリューム感を尊重してください。特に指定がない場合は、執筆物の種類に応じて一般的な長さを想定してください（例：短編小説なら数千字、ブログ記事なら1500-2000字、レポートなら数ページなど）。
- その他: ユーザーがその執筆物を通して達成したい目的（例：情報提供、娯楽、説得、記録など）を考慮してください。
"""
        finalized_requirements = (f"ユーザーの初期リクエスト: 「{original_prompt}」\n\n"
                                  f"会話全体の主要な目的: 「{initial_user_prompt_for_session}」\n\n"
                                  f"AIによる要件整理(W0):\n{defined_requirements}\n\n"
                                  f"{supplemented_info}")
        logger.info(f"ステップW0後 - 最終的な執筆要件 (補足情報込み):\n{finalized_requirements[:500].strip()}...")

        logger.info("\n執筆モード ステップW1: 構成案・プロット作成中...")
        w1_system_instruction = (
            "あなたは経験豊富な編集者、または作家、シナリオライター、リサーチャーです。\n"
            "提示された詳細な執筆要件を基に、対象となる執筆物の種類に応じた、論理的で分かりやすく、かつ読者/視聴者を引き込むような構成案を作成してください。\n"
            f"この会話全体の主要な目的は「{initial_user_prompt_for_session}」であり、ユーザーの元々のリクエストは「{original_prompt}」でした。"
        )
        w1_user_prompt = (f"以下の詳細な執筆要件に基づいて、執筆物の種類に応じた構成案（またはプロット、章立てなど）を作成してください。\n"
                          f"--- 確定した執筆要件 ---\n{finalized_requirements}\n--- 執筆要件ここまで ---\n")
        w1_res = await get_claude_response(
            request=request, prompt_text=w1_user_prompt, system_instruction=w1_system_instruction, model="claude-opus-4-20250514",
            chat_history=list(current_chat_history_for_this_turn), initial_user_prompt=initial_user_prompt_for_session, user_memories=user_memories,
        )
        steps_executed.append(w1_res)
        if w1_res.error or not w1_res.response: raise ValueError(f"執筆モード ステップW1 (構成案作成) 失敗: {w1_res.error or '応答がありませんでした。'}")
        article_outline = w1_res.response
        logger.info(f"ステップW1 - 作成された構成案/プロット (冒頭):\n{article_outline[:300].strip()}...")

        logger.info("\n執筆モード ステップW2: 初稿（ドラフト）執筆中...")
        w2_system_role_description = (
            "あなたはプロのライター、小説家、研究者、または脚本家です。\n"
            "与えられた構成案と詳細な執筆要件に基づいて、読者/視聴者にとって魅力的で分かりやすい文章コンテンツの初稿を執筆してください。\n"
            f"この会話全体の主要な目的は「{initial_user_prompt_for_session}」であり、ユーザーの元々のリクエストは「{original_prompt}」でした。"
        )
        w2_user_prompt = (f"以下の執筆要件と構成案（またはプロットなど）に基づいて、執筆物の初稿を作成してください。\n"
                          f"--- 確定した執筆要件 ---\n{finalized_requirements}\n--- 執筆要件ここまで ---\n\n"
                          f"--- 作成された構成案/プロット ---\n{article_outline}\n--- 構成案/プロットここまで ---\n")
        w2_res = await get_openai_response(
            request=request, prompt_text=w2_user_prompt, system_role_description=w2_system_role_description, model="gpt-4o",
            chat_history=list(current_chat_history_for_this_turn), initial_user_prompt=initial_user_prompt_for_session, user_memories=user_memories,
        )
        steps_executed.append(w2_res)
        if w2_res.error or not w2_res.response: raise ValueError(f"執筆モード ステップW2 (初稿執筆) 失敗: {w2_res.error or '応答がありませんでした。'}")
        draft_content = w2_res.response
        logger.info(f"ステップW2 - 作成された初稿 (冒頭):\n{draft_content[:300].strip()}...")

        logger.info("\n執筆モード ステップW3: 内容レビューと改善提案中...")
        w3_system_instruction = (
            "あなたは経験豊富なプロの編集者です。提示された執筆物の初稿、元々の執筆要件、そして構成案を注意深く読み込み、詳細なレビューと具体的な改善提案を行ってください。\n"
            f"この会話全体の主要な目的は「{initial_user_prompt_for_session}」であり、ユーザーの元々のリクエストは「{original_prompt}」でした。"
        )
        w3_user_prompt = (f"以下の執筆物の初稿について、プロの編集者として詳細なレビューと具体的な改善提案をお願いします。\n\n"
                          f"--- 元々の執筆要件 ---\n{finalized_requirements}\n--- 執筆要件ここまで ---\n\n"
                          f"--- 事前に作成された構成案/プロット ---\n{article_outline}\n--- 構成案/プロットここまで ---\n\n"
                          f"--- AIが執筆した初稿 ---\n{draft_content}\n--- 初稿ここまで ---\n\n"
                          f"上記のシステム指示に基づき、この初稿をより魅力的で質の高いコンテンツにするための、鋭い指摘と具体的な改善案をリスト形式で複数提示してください。")
        w3_res = await get_claude_response(
            request=request, prompt_text=w3_user_prompt, system_instruction=w3_system_instruction, model="claude-opus-4-20250514",
            chat_history=list(current_chat_history_for_this_turn), initial_user_prompt=initial_user_prompt_for_session, user_memories=user_memories,
        )
        steps_executed.append(w3_res)
        if w3_res.error or not w3_res.response: raise ValueError(f"執筆モード ステップW3 (レビュー) 失敗: {w3_res.error or '応答がありませんでした。'}")
        review_and_suggestions = w3_res.response
        logger.info(f"ステップW3 - レビューと改善提案 (冒頭):\n{review_and_suggestions[:300].strip()}...")

        logger.info("\n執筆モード ステップW4: 推敲・リライト中...")
        w4_preamble_cohere = (
            "あなたはプロの編集者兼ライターです。提示された執筆物の初稿と、それに対する詳細なレビューおよび改善提案を深く理解し、レビューでの指摘事項を的確に反映させてコンテンツ全体を全面的に推敲・リライトしてください。\n"
            f"この会話全体の主要な目的は「{initial_user_prompt_for_session}」であり、ユーザーの元々のリクエストは「{original_prompt}」でした。\n"
        )
        w4_user_prompt = (f"以下の執筆物の初稿と、それに対するレビューおよび改善提案があります。\n"
                          f"これらを基に、初稿を全面的に推敲・リライトし、より質の高いコンテンツ（第2稿）を作成してください。\n\n"
                          f"--- AIが執筆した初稿 ---\n{draft_content}\n--- 初稿ここまで ---\n\n"
                          f"--- 上記初稿へのレビューと改善提案 ---\n{review_and_suggestions}\n--- レビューと改善提案ここまで ---\n\n")
        w4_res = await get_cohere_response(
            request=request, prompt_text=w4_user_prompt, preamble=w4_preamble_cohere, model="command-a-03-2025",
            chat_history=list(current_chat_history_for_this_turn), initial_user_prompt=initial_user_prompt_for_session, user_memories=user_memories,
        )
        steps_executed.append(w4_res)
        if w4_res.error or not w4_res.response: raise ValueError(f"執筆モード ステップW4 (推敲・リライト) 失敗: {w4_res.error or '応答がありませんでした。'}")
        revised_draft_content = w4_res.response
        logger.info(f"ステップW4 - 推敲・リライトされた第2稿 (冒頭):\n{revised_draft_content[:300].strip()}...")

        logger.info("\n執筆モード ステップW5: 最終校正と仕上げ中...")
        w5_system_instruction_parts = [
            "あなたは最高の編集長兼最終ライターです。与えられた執筆物（第2稿）を、以下の指示に従って完璧な最終版に仕上げてください。",
            "最終出力は小説・物語本文のみとし、章立て案や解説、ストーリー案内、コメントなど本文以外の要素は一切含めないでください。出力が途中で止まった場合は、続きの物語本文のみを書き足してください。",
            f"\nこの会話全体の主要な目的は「{initial_user_prompt_for_session}」であり、ユーザーの元々のリクエストは「{original_prompt}」であったことを常に念頭に置いてください。"
        ]
        w5_final_system_instruction = "\n\n".join(w5_system_instruction_parts)
        w5_user_prompt = (f"以下の執筆物（第2稿）を、上記のシステム指示に従って最終校正し、指定された口調（もしあれば）で最高の形に仕上げてください。\n"
                          "出力は物語本文のみとし、案内文やコメント、章立て解説を一切含めないでください。続きが必要な場合は本文だけを書き足してください。\n\n"
                          f"--- 推敲・リライトされた執筆物（第2稿） ---\n{revised_draft_content}\n--- 第2稿ここまで ---\n\n")
        w5_res = await get_gemini_response(
            request=request, prompt_text=w5_user_prompt, system_instruction=w5_final_system_instruction, model_name="gemini-2.5-pro-preview-05-06",
            chat_history=list(current_chat_history_for_this_turn), initial_user_prompt=initial_user_prompt_for_session, user_memories=user_memories,
        )
        steps_executed.append(w5_res)
        if w5_res.error or not w5_res.response: raise ValueError(f"執筆モード ステップW5 (最終校正・仕上げ) 失敗: {w5_res.error or '応答がありませんでした。'}")
        final_article_content = w5_res.response
        logger.info(f"ステップW5 - 完成版コンテンツ (冒頭):\n{final_article_content[:300].strip()}...")

        final_step_response = schemas.IndividualAIResponse(source="Writing Mode (W5 - Final Content)", response=final_article_content)
        response_shell.step7_final_answer_v2_openai = final_step_response
        response_shell.writing_mode_details = steps_executed
        logger.info("--- 執筆特化モード (全ステップ完了) 終了 ---")
    except ValueError as e:
        error_message = f"執筆特化モードの処理中にエラー: {str(e)}"
        logger.info(error_message)
        response_shell.overall_error = error_message
        response_shell.writing_mode_details = steps_executed
        if not response_shell.step7_final_answer_v2_openai:
            response_shell.step7_final_answer_v2_openai = schemas.IndividualAIResponse(
                source="Writing Mode Error Step", error=str(e), response=f"申し訳ありません、処理中にエラーが発生しました。\nエラー内容: {str(e)}",
            )
    return response_shell


# --- 超長文執筆モード ---
async def run_ultra_writing_mode_flow(
    original_prompt: str,
    response_shell: schemas.CollaborativeResponseV2,
    chat_history_for_ai: List[Dict[str, str]],
    initial_user_prompt_for_session: Optional[str],
    user_memories: Optional[List[schemas.UserMemoryResponse]] = None,
    desired_char_count: Optional[int] = None,
    request: Request = None,
) -> schemas.CollaborativeResponseV2:
    logger.info("\n--- 超長文執筆モード開始 ---")
    # Ensure `request` is used for AI helper calls within this function.
    steps_executed: List[schemas.IndividualAIResponse] = []
    logger.info(f"Ultra Writing Mode: メモリ: {len(user_memories) if user_memories else 0}件")
    current_context_history = list(chat_history_for_ai)
    current_context_history.append({"role": "user", "content": original_prompt})

    try:
        outline_res = await get_openai_response(
            request=request, # Pass request
            prompt_text=f"次の内容で章立てを提案してください:\n{original_prompt}",
            system_role_description="Long Form Outline Generator",
            chat_history=current_context_history,
            initial_user_prompt=initial_user_prompt_for_session,
            user_memories=user_memories,
        )
        steps_executed.append(outline_res)
        if outline_res.error or not outline_res.response: raise ValueError("構成案生成に失敗しました")

        chapters = [c.strip() for c in outline_res.response.splitlines() if c.strip()]
        final_text = ""
        for ch_idx, ch_title in enumerate(chapters):
            logger.info(f"  章 {ch_idx+1}/{len(chapters)} 「{ch_title}」を執筆中...")
            chapter_prompt = (f"以下の章「{ch_title}」について、詳細な本文を執筆してください。")
            if final_text: chapter_prompt += f"\n\nこれまでのあらすじや主要な流れを簡単に振り返ると、『{final_text[-500:]}...』といった内容でした。\nこれを踏まえて執筆を続けてください。"
            chapter_res = await get_openai_response(
                request=request, # Pass request
                prompt_text=chapter_prompt,
                system_role_description="Chapter Writer: 本文のみを出力し、章タイトルや案内文は含めないでください。詳細かつ具体的に記述してください。",
                chat_history=current_context_history,
                initial_user_prompt=initial_user_prompt_for_session,
                user_memories=user_memories,
            )
            steps_executed.append(chapter_res)
            if chapter_res.response:
                final_text += f"\n## {ch_title}\n\n{chapter_res.response.strip()}\n"
                current_context_history.append({"role": "assistant", "content": chapter_res.response.strip()})
                if ch_idx + 1 < len(chapters): current_context_history.append({"role": "user", "content": f"ありがとうございます。次の章『{chapters[ch_idx+1]}』に進んでください。"})
                else: current_context_history.append({"role": "user", "content": "ありがとうございます。これで全ての章が完了しました。"})

        if desired_char_count and isinstance(desired_char_count, int) and desired_char_count > 0:
            loop_count = 0
            max_expansion_loops = 5
            while len(final_text) < desired_char_count and loop_count < max_expansion_loops:
                loop_count += 1
                remaining_chars = desired_char_count - len(final_text)
                logger.info(f"  文字数調整ループ {loop_count}: 残り約{remaining_chars}文字...")
                expansion_prompt = (f"現在の文章は以下の通りです。\n{final_text[-1000:]}...\n\n"
                                    f"この文章全体の内容をさらに詳細に、具体的に、物語であれば描写を豊かに、説明文であれば具体例や補足情報を加えて拡張してください。")
                add_res = await get_openai_response(
                    request=request, # Pass request
                    prompt_text=expansion_prompt,
                    system_role_description="Expansion Writer: 提供された文章を自然な形で加筆・拡張し、詳細を豊かにしてください。本文のみを続けて出力してください。コメントや案内は禁止です。",
                    chat_history=current_context_history,
                    initial_user_prompt=initial_user_prompt_for_session,
                    user_memories=user_memories,
                )
                steps_executed.append(add_res)
                if not add_res.response or not add_res.response.strip(): break
                final_text += f"\n\n{add_res.response.strip()}\n"
                current_context_history.append({"role": "assistant", "content": add_res.response.strip()})
                if len(final_text) < desired_char_count: current_context_history.append({"role": "user", "content": "ありがとうございます。さらに内容を拡張してください。"})

        response_shell.step7_final_answer_v2_openai = schemas.IndividualAIResponse(source="Ultra LongWriting Final", response=final_text)
        response_shell.ultra_writing_mode_details = steps_executed
        logger.info("--- 超長文執筆モード終了 ---")
    except Exception as e:
        error_msg = f"超長文執筆モードの処理中にエラー: {str(e)}"
        logger.info(error_msg)
        response_shell.overall_error = error_msg
        response_shell.ultra_writing_mode_details = steps_executed
        if not response_shell.step7_final_answer_v2_openai:
            response_shell.step7_final_answer_v2_openai = schemas.IndividualAIResponse(
                source="Ultra LongWriting Error", error=str(e), response=f"申し訳ありません、処理中にエラーが発生しました。\nエラー内容: {str(e)}",
            )
    return response_shell


async def run_fast_chat_mode_flow(
    original_prompt: str,
    response_shell: schemas.CollaborativeResponseV2,
    chat_history_for_ai: List[Dict[str, str]],
    initial_user_prompt_for_session: Optional[str],
    user_memories: Optional[List[schemas.UserMemoryResponse]] = None,
    model: str = "gpt-4o",
    request: Request = None,
) -> schemas.CollaborativeResponseV2:
    logger.info("\n--- 高速チャットモード開始 ---")
    logger.info(
        f"Fast Chat Mode: メモリ: {len(user_memories) if user_memories else 0}件, 履歴件数(AIへ): {len(chat_history_for_ai)}"
    )
    res = await get_openai_response(
        request=request, # Pass request object
        prompt_text=original_prompt,
        system_role_description="Fast Chat Mode",
        model=model,
        chat_history=list(chat_history_for_ai),
        initial_user_prompt=initial_user_prompt_for_session,
        user_memories=user_memories,
    )
    response_shell.step7_final_answer_v2_openai = res
    logger.info("--- 高速チャットモード終了 ---")
    return response_shell

    # サーバー起動コマンド (開発時): uvicorn main:app --reload
